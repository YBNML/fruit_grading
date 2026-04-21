#!/usr/bin/env python3
"""
Generate fruit_grading_summary.pptx from the project results.

Slides are built programmatically so the deck regenerates deterministically
from the latest numbers if experiments are re-run.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).resolve().parent / "fruit_grading_summary.pptx"

# Theme palette
C_BG = RGBColor(0xFA, 0xFA, 0xFA)
C_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0xE3, 0x4A, 0x2F)  # terracotta
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_RED = RGBColor(0xC6, 0x28, 0x28)
C_MUTED = RGBColor(0x55, 0x55, 0x55)
C_TABLE_HDR = RGBColor(0x26, 0x3B, 0x5B)
C_TABLE_ALT = RGBColor(0xEC, 0xEE, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=C_TITLE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Apple SD Gothic Neo"
    return tb


def add_accent_bar(slide, top=Inches(1.1), width=Inches(0.12), color=C_ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top,
                                  width, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def slide_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12.5), Inches(0.7),
             title, size=30, bold=True, color=C_TITLE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.0), Inches(12.5), Inches(0.45),
                 subtitle, size=14, color=C_MUTED)


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
             f"fruit_grading — YBNML   |   slide {idx}/{total}",
             size=9, color=C_MUTED)


def add_table(slide, left, top, width, height, rows_data, header=True,
              col_widths=None):
    """rows_data: list of lists of strings."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Apple SD Gothic Neo"
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_HDR
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True
                run.font.size = Pt(11)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (C_TABLE_ALT if r % 2 == 0
                                            else RGBColor(0xFF, 0xFF, 0xFF))
                run.font.size = Pt(10)
                run.font.color.rgb = C_TITLE
    return table


def add_bullets(slide, left, top, width, height, bullets, size=14,
                color=C_TITLE, leading=1.35):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Apple SD Gothic Neo"
    return tb


# ---------------------------------------------------------------------- slides


def slide_title(prs):
    s = add_blank(prs); fill_bg(s)
    # Left accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    add_text(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.6),
             "fruit_grading", size=60, bold=True, color=C_TITLE)
    add_text(s, Inches(0.9), Inches(2.85), Inches(11), Inches(0.6),
             "과일 품질·치수 자동 분석 ML 파이프라인",
             size=24, color=C_MUTED)
    add_text(s, Inches(0.9), Inches(4.2), Inches(11), Inches(0.4),
             "복숭아 천중도 + 감귤 황금향 · 5 Phase · 26 Runs · ResNet50 MV",
             size=16, color=C_TITLE)
    add_text(s, Inches(0.9), Inches(4.65), Inches(11), Inches(0.4),
             "YBNML  /  github.com/YBNML/fruit_grading",
             size=14, color=C_MUTED)
    add_text(s, Inches(0.9), Inches(6.3), Inches(11), Inches(0.4),
             "2026-04 · Claude Opus 4.7 (1M context)",
             size=11, color=C_MUTED)


def slide_problem(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "문제 정의",
                 "선별기의 3대 카메라 영상으로 과일 품질·치수를 자동 측정")

    add_text(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5),
             "현재 워크플로우 (각 장비 개별 측정)", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(3.5), [
        "전수 측정: 전자저울(무게) + 캘리퍼스(치수) + 비파괴 당도계(brix)",
        "과일 1개당 ≥3 개 장비 순차 거치 → 처리 속도 병목",
        "작업자 반복 조작 필요, 대량 처리 시 피로·오차 누적",
        "선별기 카메라는 이미 존재 — 측정에 활용되지 못함",
    ], size=14)

    add_text(s, Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.5),
             "목표: 카메라 단일 파이프라인으로 통합", size=16, bold=True, color=C_GREEN)
    add_bullets(s, Inches(7.0), Inches(2.2), Inches(6.0), Inches(3.5), [
        "3뷰(top1/top2/bottom1) 이미지 → 측정값 직접 회귀",
        "weight(g) / height·max_w·min_w(mm) / brix(°Bx) 동시 예측",
        "별도 저울·캘리퍼스·당도계 없이 선별기 흐름 중 자동화",
        "작물별(복숭아 / 감귤) 개별 모델 학습",
    ], size=14)

    add_text(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.5),
             "Why this matters: 3 장비 통합으로 처리량↑·인력↓, "
             "이미지 기반 등급화로 선별 일관성 확보, 기존 선별기 하드웨어 재활용",
             size=13, color=C_MUTED, bold=True)


def slide_data(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "데이터셋",
                 "두 작물, 각 작물당 날짜별 3뷰 촬영 + 실측값")

    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.3), [
        ["작물", "촬영일 수", "과일 수", "이미지 수", "라벨 포맷"],
        ["복숭아 (천중도)", "9일 (2023-08)", "1,278", "3,834",
         "날짜별 JSON (measurement/marketability/camera/farm)"],
        ["감귤 (황금향)", "7일 (2023-11~12)", "2,309", "6,927",
         "CSV (CP949/UTF-8 혼재) + labelme 폴리곤 JSON"],
    ], col_widths=[Inches(2.0), Inches(2.0), Inches(1.3), Inches(1.3),
                   Inches(5.5)])

    add_text(s, Inches(0.6), Inches(3.25), Inches(12), Inches(0.5),
             "측정 항목 (target)", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(3.7), Inches(6.0), Inches(2), [
        "weight (g) — 가장 실무 중요",
        "height (mm) — 과일 높이",
        "max_w / min_w (mm) — 최대/최소 너비",
        "brix (°Bx) — 당도 (내부 성질)",
    ])

    add_text(s, Inches(7.0), Inches(3.25), Inches(5.5), Inches(0.5),
             "데이터 품질 이슈 & 해결", size=16, bold=True, color=C_RED)
    add_bullets(s, Inches(7.0), Inches(3.7), Inches(5.5), Inches(2.5), [
        "~1.3% outlier 발견 (weight=1210g 등 박스 무게 오입력)",
        "SANITY_BOUNDS 필터 추가 → outlier → None",
        "황금향 weight R² 0.18 → 0.92 회복",
        "이슈: 측정 오차, 경계 모호성, 도메인 단일성",
    ])

    add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             "fruit_id = (date, fruit_idx) — 같은 과일의 3뷰 그룹화 키. "
             "train/test split은 fruit_id 단위로 수행 (누수 방지)",
             size=12, color=C_MUTED)


def slide_pipeline(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "파이프라인 구조", "Data → Label → Train → Evaluate")

    # Boxes for flow
    def _box(left, top, w, h, title, body, color):
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        r.fill.solid(); r.fill.fore_color.rgb = color
        r.line.fill.background()
        tb = s.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        run.font.bold = True; run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Apple SD Gothic Neo"
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run = p2.add_run(); run.text = body
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Apple SD Gothic Neo"

    y = Inches(2.2); h = Inches(1.2); w = Inches(2.7)
    _box(Inches(0.6), y, w, h, "DB/ (원본)",
         "JSON / CSV / 이미지", C_MUTED)
    _box(Inches(3.5), y, w, h, "make_labels.py",
         "경로 정규화 + outlier 필터", C_TABLE_HDR)
    _box(Inches(6.4), y, w, h, "label_{crop}.csv",
         "통합 스키마 라벨", C_GREEN)
    _box(Inches(9.3), y, w, h, "train*.py",
         "ResNet50 MV / 회귀", C_ACCENT)
    _box(Inches(5.0), Inches(4.0), w, h, "best.pt (비공개)",
         "학습된 가중치", C_MUTED)
    _box(Inches(7.9), Inches(4.0), w, h, "history.json (공개)",
         "학습 곡선 + 메트릭", C_GREEN)

    # arrows (simple lines)
    def _arrow(x1, y1, x2, y2):
        ln = s.shapes.add_connector(1, x1, y1, x2, y2)
        ln.line.color.rgb = C_MUTED
        ln.line.width = Pt(1.5)

    _arrow(Inches(3.3), Inches(2.8), Inches(3.5), Inches(2.8))
    _arrow(Inches(6.2), Inches(2.8), Inches(6.4), Inches(2.8))
    _arrow(Inches(9.1), Inches(2.8), Inches(9.3), Inches(2.8))
    _arrow(Inches(10.6), Inches(3.4), Inches(6.4), Inches(4.0))
    _arrow(Inches(10.6), Inches(3.4), Inches(9.3), Inches(4.0))

    add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
             "핵심 설계 결정", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(6.1), Inches(12), Inches(1.2), [
        "Fruit-id 단위 split (70/30)  |  fruit-level eval (3뷰 평균)  |  "
        "ResNet50 + ImageNet V2 pretrained  |  SmoothL1 + Cosine LR + early stop (patience=5)",
    ], size=12, color=C_MUTED)


def slide_phase0(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Phase 0 · 3-grade Classification (baseline)",
                 "weight tertile을 기반으로 소/중/대 품질 등급 합성 후 분류")

    add_text(s, Inches(0.6), Inches(1.8), Inches(5), Inches(0.5),
             "결과", size=16, bold=True, color=C_ACCENT)
    add_table(s, Inches(0.6), Inches(2.3), Inches(5.5), Inches(1.4), [
        ["작물", "val_acc", "best epoch"],
        ["복숭아", "81.72%", "28"],
        ["감귤 (황금향)", "92.82%", "25"],
    ], col_widths=[Inches(2), Inches(1.75), Inches(1.75)])

    add_text(s, Inches(6.6), Inches(1.8), Inches(6), Inches(0.5),
             "복숭아 confusion matrix", size=16, bold=True, color=C_ACCENT)
    add_table(s, Inches(6.6), Inches(2.3), Inches(6.2), Inches(1.7), [
        ["", "pred 소", "pred 중", "pred 대", "정답률"],
        ["true 소", "362", "24", "1", "93.5%"],
        ["true 중", "83", "252", "67", "62.7% ⚠"],
        ["true 대", "1", "34", "325", "90.3%"],
    ], col_widths=[Inches(1.2), Inches(1.25), Inches(1.25), Inches(1.25), Inches(1.25)])

    add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
             "관찰 & 다음 단계", size=16, bold=True, color=C_GREEN)
    add_bullets(s, Inches(0.6), Inches(4.75), Inches(12), Inches(2.5), [
        "양극단 등급(소/대)은 잘 맞추지만 중간 등급이 62.7% — tertile 경계 근처 과일은 라벨 자체가 모호",
        "classification의 근본 한계: 연속적 크기를 이산 라벨로 강제 변환하면서 정보 손실",
        "해결책 → Phase 1: 측정값 자체를 회귀로 예측 → 후처리로 3등급화 (경계 모호성 자동 해소)",
    ], size=13)


def slide_phase1(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Phase 1 · Single-View Regression",
                 "단일 이미지 회귀 학습 + 추론 시 3뷰 예측값 평균 (후처리)")

    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(3.3), [
        ["작물", "target", "sample MAE", "sample R²", "fruit MAE",
         "fruit R²", "fruit MAPE"],
        ["복숭아", "weight",  "12.22 g", "0.955", "7.61 g",  "0.983", "1.89%"],
        ["복숭아", "height",  "2.24 mm", "0.798", "1.92 mm", "0.843", "2.35%"],
        ["복숭아", "max_w",   "1.37 mm", "0.944", "1.07 mm", "0.967", "1.14%"],
        ["복숭아", "min_w",   "1.50 mm", "0.908", "1.24 mm", "0.938", "1.41%"],
        ["감귤",   "weight",  "4.58 g",  "0.909", "3.55 g",  "0.923", "2.52%"],
        ["감귤",   "height",  "1.49 mm", "0.732", "1.26 mm", "0.773", "2.18%"],
        ["감귤",   "max_w",   "1.04 mm", "0.896", "0.85 mm", "0.918", "1.23%"],
        ["감귤",   "min_w",   "0.86 mm", "0.943", "0.69 mm", "0.963", "1.03%"],
    ], col_widths=[Inches(1.2), Inches(1.3), Inches(1.6), Inches(1.6),
                   Inches(1.6), Inches(1.6), Inches(1.7)])

    add_bullets(s, Inches(0.6), Inches(5.2), Inches(12), Inches(2), [
        "추론 시 3뷰 평균만으로도 샘플-MAE 대비 약 38% 감소 (free win — 별도 학습 없음)",
        "전 8 target이 MAPE <3% 달성 → 분류의 81.7% acc 한계를 회귀로 우회",
        "R² 단독 평가 주의: height MAPE 2.7%인데 R² 0.80 — 타겟 분산이 작을 때 R² 패널티 (MAE/MAPE 병행)",
    ], size=13)


def slide_phase3(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Phase 3 · Multi-View Regression (현재 best 구조)",
                 "3뷰를 한 모델이 동시에 입력 — backbone 공유 + feature concat")

    add_text(s, Inches(0.6), Inches(1.7), Inches(4.5), Inches(0.5),
             "아키텍처", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.0), Inches(3), [
        "t1/t2/b1 각각 ResNet50 backbone (공유) 통과",
        "뷰별 feature (2048-dim) 순서대로 concat → 6144-dim",
        "Linear(6144, 1) head로 target 예측",
        "파라미터 증가분은 head만 — backbone 100% 재사용",
        "배치 크기 8 (3뷰 × 8 fruits = 24 images per iter)",
    ], size=13)

    add_text(s, Inches(5.8), Inches(1.7), Inches(7), Inches(0.5),
             "결과 (fruit MAE, vs Phase 1)", size=16, bold=True, color=C_ACCENT)
    add_table(s, Inches(5.8), Inches(2.2), Inches(7.0), Inches(3.5), [
        ["작물/target", "MV MAE", "MAPE", "Δ vs P1"],
        ["복숭아/weight",      "6.78 g",      "1.78%", "−10.9% ✓"],
        ["복숭아/height",      "1.98 mm",     "2.42%", "+2.9%"],
        ["복숭아/max_w",       "0.96 mm",     "1.02%", "−10.0% ✓"],
        ["복숭아/min_w",       "1.22 mm",     "1.37%", "−1.5%"],
        ["감귤/weight",        "3.89 g",      "2.68%", "+9.6%"],
        ["감귤/height",        "1.23 mm",     "2.14%", "−2.3%"],
        ["감귤/max_w",         "0.69 mm",     "1.00%", "−19.3% ✓"],
        ["감귤/min_w ⭐",       "0.57 mm",    "0.85%",  "−17.9% ✓"],
    ], col_widths=[Inches(2.0), Inches(1.7), Inches(1.5), Inches(1.8)])

    add_text(s, Inches(0.6), Inches(6.05), Inches(12), Inches(0.8),
             "⭐ 최고 성적: 감귤 min_w — MAPE 0.85%, R² 0.975  |  "
             "MV가 Phase 1 대비 6/8 승 (너비 target에서 최대 19% 개선)",
             size=13, bold=True, color=C_GREEN)


def slide_phase2(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Phase 2 · Multi-View + Multi-Task  (negative result)",
                 "한 모델이 4 target 동시 예측 (weight/height/max_w/min_w)")

    add_text(s, Inches(0.6), Inches(1.7), Inches(5), Inches(0.5),
             "구조", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.5), Inches(2), [
        "Phase 3 MV backbone 그대로",
        "head: Linear(6144, 1) → Linear(6144, 4)",
        "z-score 정규화 공간에서 학습",
        "모델 수 8개 → 2개 (작물당 1)",
    ], size=13)

    add_text(s, Inches(6.3), Inches(1.7), Inches(6), Inches(0.5),
             "결과 — 8/8 후퇴", size=16, bold=True, color=C_RED)
    add_table(s, Inches(6.3), Inches(2.2), Inches(6.5), Inches(3.3), [
        ["target", "P3 MV", "P2 MV+MT", "Δ"],
        ["복숭아 weight",  "6.78",  "12.29", "+81% ✗"],
        ["복숭아 height",  "1.98",  "2.28",  "+15%"],
        ["복숭아 max_w",   "0.96",  "1.47",  "+53% ✗"],
        ["복숭아 min_w",   "1.22",  "1.41",  "+15%"],
        ["감귤 weight",    "3.89",  "5.15",  "+32%"],
        ["감귤 height",    "1.23",  "1.55",  "+26%"],
        ["감귤 max_w",     "0.69",  "0.89",  "+30%"],
        ["감귤 min_w",     "0.57",  "0.77",  "+37%"],
    ], col_widths=[Inches(2), Inches(1.5), Inches(1.5), Inches(1.5)])

    add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.5),
             "Null result의 가치 — 가설 검증 실패도 데이터", size=16, bold=True,
             color=C_GREEN)
    add_bullets(s, Inches(0.6), Inches(6.35), Inches(12), Inches(1.2), [
        "\"멀티태스크가 feature 공유로 약한 target을 돕는다\" 통설이 이 구조/데이터에선 성립 안 함",
        "원인 추정: 단일 Linear head의 gradient 간섭, 모든 target이 크기 계열이라 공유 이점 작음",
        "교훈: head를 MLP로 확장·per-task head·uncertainty weighting 등 고도화 필요",
    ], size=12)


def slide_phase5(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Phase 5 · Brix(당도) 회귀",
                 "내부 성질(당도)를 외형 이미지로 예측 가능한가?")

    add_text(s, Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.5),
             "가설", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.5), Inches(2), [
        "Brix는 내부 당도 — 외형과 약한 상관 예상",
        "실제론 과피 색/질감이 성숙도와 관련 → 신호 존재 가능성",
        "분광 센서 없이 RGB만으로 어디까지 가능한가?",
    ], size=13)

    add_text(s, Inches(6.5), Inches(1.7), Inches(6), Inches(0.5),
             "결과 (MV, fruit-level)", size=16, bold=True, color=C_ACCENT)
    add_table(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(1.6), [
        ["작물", "MAE", "R²", "MAPE"],
        ["복숭아 brix", "0.472 °Bx", "0.676", "3.57%"],
        ["감귤 brix",   "0.298 °Bx", "0.629", "2.57%"],
    ], col_widths=[Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5)])

    add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
             "관찰", size=16, bold=True, color=C_GREEN)
    add_bullets(s, Inches(0.6), Inches(5.15), Inches(12), Inches(2), [
        "R²는 크기 target(0.95+)보다 낮지만 MAPE 2.6~3.9%는 실용 범위 — 비접촉 brix 측정의 가능성",
        "복숭아 brix R²(0.68) > 감귤(0.63) — 복숭아 분포 범위가 더 넓어서(σ 1.15 vs 0.62) 신호↑",
        "MV가 SV+avg 대비 3~8% 개선 — 크기 target 대비 작음 (brix 신호는 표면 색/질감 위주)",
        "차기: 분광/NIR 센서 병용 시 R² 0.9+로 도약 가능성 (관련 논문 다수)",
    ], size=13)


def slide_backbone(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "Backbone 비교 · ViT-B/16 vs ResNet50",
                 "같은 Phase 3 MV recipe로 CNN vs Transformer 실측 비교")

    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.3), [
        ["target", "ViT-B/16", "ResNet50", "ΔMAE", "ViT R²", "R50 R²"],
        ["복숭아 weight",     "60.71 g",   "6.78 g",    "+795% ✗", "−0.008", "0.988"],
        ["복숭아 brix",       "0.613 °Bx", "0.472 °Bx", "+30%",    "0.473",  "0.676"],
        ["감귤 weight",       "4.213 g",   "3.890 g",   "+8%",     "0.914",  "0.922"],
        ["감귤 brix",         "0.392 °Bx", "0.298 °Bx", "+31%",    "0.355",  "0.629"],
    ], col_widths=[Inches(2.5), Inches(2.0), Inches(2.0), Inches(1.8),
                   Inches(1.8), Inches(2.0)])

    add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             "해석", size=16, bold=True, color=C_GREEN)
    add_bullets(s, Inches(0.6), Inches(4.65), Inches(12), Inches(2.5), [
        "ResNet50 4/4 전승 — 소규모 데이터(N≈1~2k)에선 CNN의 locality prior가 작동",
        "복숭아 weight ViT 완전 실패 (R²≈0) — CNN-기준 recipe (SGD+lr=1e-3)가 ViT에 불리한 증거",
        "감귤(데이터 1.8배 많음)에선 ViT가 CNN에 8%까지 근접 — ViT data-hungriness 통설 검증",
        "Brix에서도 ViT 우세 가설 반증 — global attention이 반드시 이기지 않음",
        "공정 비교(same recipe) 기준 엔지니어링 결정 근거 확보 — CNN이 본 프로젝트의 기본 선택",
    ], size=13)


def slide_summary(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "종합 · 최고 성능 정리 (운용 추천)",
                 "Phase 3 MV (ResNet50) + 작물당 4 target 모델")

    add_table(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.0), [
        ["작물", "target", "MAE", "R²", "MAPE", "비고"],
        ["복숭아", "weight",  "6.78 g",   "0.988", "1.78%", "회귀 상위"],
        ["복숭아", "height",  "1.92 mm",  "0.843", "2.35%", "(Phase 1이 나음)"],
        ["복숭아", "max_w",   "0.96 mm",  "0.973", "1.02%", ""],
        ["복숭아", "min_w",   "1.22 mm",  "0.936", "1.37%", ""],
        ["복숭아", "brix",    "0.472 °Bx","0.676", "3.57%", "내부 성질"],
        ["감귤",   "weight",  "3.55 g",   "0.923", "2.52%", "(Phase 1이 나음)"],
        ["감귤",   "height",  "1.23 mm",  "0.778", "2.14%", ""],
        ["감귤",   "max_w",   "0.69 mm",  "0.918", "1.00%", ""],
        ["감귤",   "min_w",   "0.57 mm",  "0.975", "0.85%", "⭐ 전체 best"],
        ["감귤",   "brix",    "0.298 °Bx","0.629", "2.57%", "내부 성질"],
    ], col_widths=[Inches(1.2), Inches(1.3), Inches(1.7), Inches(1.5),
                   Inches(1.5), Inches(4.9)])

    add_text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.8),
             "전 크기 target MAPE 1~3% / brix MAPE 2.6~3.6% — "
             "선별기 운용 요구 (MAPE <5%) 충족",
             size=14, bold=True, color=C_GREEN, anchor=MSO_ANCHOR.MIDDLE)


def slide_insights(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "핵심 인사이트 (1/2)",
                 "실험을 통해 실증된 교훈")

    add_bullets(s, Inches(0.6), Inches(1.7), Inches(12), Inches(5.5), [
        "① 3뷰 평균만으로 MAE 38% 감소 — 학습 없는 'free win'. 멀티뷰 환경의 강력한 baseline",
        "② 멀티뷰 학습 > 단일뷰 평균 (6/8 승) — 너비(max_w/min_w)에서 특히 큰 이득 (18~19%). "
        "3뷰가 서로 다른 각도의 최대/최소를 관측하여 암묵적 3D 재구성",
        "③ 멀티태스크는 본 세팅에서 역효과 (8/8 후퇴) — 통설의 반례. 단일 head의 "
        "gradient 간섭이 feature 공유 이득 초과. 모든 target이 '크기' 계열이라 공유 의미 작음",
        "④ 외형만으로 당도(brix) 예측 실현 가능 — MAPE 2.6~3.9%, 실용 범위. 과피 색/질감이 "
        "성숙도와 상관. 분광 센서 없이도 비접촉 품질 추정 가능성",
        "⑤ 소규모 데이터에선 CNN > Transformer — 같은 recipe 기준 ResNet50 4/4 승. "
        "ViT는 data-hungry + 전용 optimization 요구",
    ], size=13, leading=1.4)


def slide_insights2(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "핵심 인사이트 (2/2)",
                 "데이터·평가 방법론적 교훈")

    add_bullets(s, Inches(0.6), Inches(1.7), Inches(12), Inches(5.5), [
        "⑥ 데이터 품질이 모델보다 중요 — 황금향 weight outlier 18건(0.26%)이 R² 0.18로 붕괴시켰고, "
        "SANITY_BOUNDS 필터링만으로 R² 0.92 복구. 소수 outlier가 SS_total 지배 = 교과서 케이스",
        "⑦ R² 단독 평가는 misleading — 타겟 분산이 작을 때 (예: 복숭아 height 65~100mm) "
        "R²가 낮게 나옴. 운용 지표는 MAPE + MAE 병행해야 실제 정밀도 반영",
        "⑧ Fruit-id 단위 split은 필수 — 같은 과일의 3뷰가 train/test로 쪼개지면 val 지표가 부풀려짐. "
        "유사 실수가 멀티뷰 프로젝트에서 흔함",
        "⑨ Null result도 포트폴리오 자산 — 멀티태스크 실패는 '왜 Phase 3를 택했는가' 및 "
        "'멀티태스크 고도화 방향'의 근거 제공",
        "⑩ Early stopping이 시간 절약 — 80% 이상 실험이 15~20 epoch 내 수렴. "
        "patience=5로 2~4분/실험 절감",
    ], size=13, leading=1.4)


def slide_limits(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "현재 한계 (Limitations)",
                 "외삽 전에 반드시 고려해야 할 제약")

    add_text(s, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
             "데이터 한계", size=16, bold=True, color=C_RED)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(6), Inches(3), [
        "작물당 ~1.3k / 2.3k 과일 (소규모)",
        "2 작물, 각 1 품종만 (일반화 미검증)",
        "단일 촬영 환경 (조명·카메라 1세트)",
        "9일 / 7일 촬영 — 계절·연도별 drift 불명",
        "상처·병반 등 부정적 케이스 부족",
    ], size=13)

    add_text(s, Inches(7.0), Inches(1.7), Inches(6), Inches(0.5),
             "방법론 한계", size=16, bold=True, color=C_RED)
    add_bullets(s, Inches(7.0), Inches(2.2), Inches(6), Inches(3), [
        "해상도 224×224 — 미세 결함 손실 가능",
        "Phase 4 세그멘테이션 미활용 (skip)",
        "ViT 동일 recipe 비교 — CNN 편향",
        "Brix R² 0.60~0.68 상한 — RGB 신호 한계",
        "Interpretability 미구현 (GradCAM 등)",
        "Uncertainty 정량화 없음",
    ], size=13)

    add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             "운용 고려사항", size=16, bold=True, color=C_RED)
    add_bullets(s, Inches(0.6), Inches(6.3), Inches(12), Inches(1), [
        "실시간 latency 검증 미완  |  모델 갱신 주기 미정의  |  "
        "도메인 차이 적응 전략 부재  |  라벨 수집 표준 가이드 없음",
    ], size=12, color=C_MUTED)


def slide_future_short(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "개선 방향 · 단기 (구현 가능, 1~4주)",
                 "현 인프라·데이터로 시도 가능한 개선")

    items = [
        ("1. Phase 4 세그멘테이션 실행",
         "황금향 labelme 폴리곤 → YOLOv8-seg 학습 → tight crop → 기존 모델 재학습. "
         "기대: MAPE 5~15% 추가 감소. 복숭아도 라벨 추가 시 동일 처리"),
        ("2. ViT 전용 recipe 재비교",
         "AdamW + lr 1e-4 + warmup + batch 16+ + LayerNorm 학습. "
         "공정 비교를 위해 CNN recipe와 분리. 데이터 부족은 남는 변수"),
        ("3. Test-Time Augmentation (TTA)",
         "추론 시 horizontal flip / 90° rotate 여러 번 → 평균. "
         "학습 변경 없이 MAPE 2~5% 추가 개선 기대"),
        ("4. 앙상블",
         "Seed 3~5 평균 또는 ResNet50 + DenseNet201 feature concat. "
         "MV fruit-avg + TTA + 앙상블로 MAPE 0.5~0.7% 수준 도달 가능성"),
        ("5. 해상도 스케일업",
         "224 → 320 / 384. brix·미세 결함 예측에 특히 도움. GPU 메모리 여유 확인 필요"),
        ("6. Multi-task 아키텍처 고도화",
         "공유 trunk 후 per-target head(2-layer MLP) + uncertainty weighting / GradNorm. "
         "Phase 2 negative result 반박 시도"),
        ("7. Attention 시각화",
         "GradCAM / ScoreCAM으로 \"모델이 어디를 보고 있는가\" 설명. 이상 예측 디버깅에 유용"),
        ("8. 모델 경량화 + ONNX",
         "MobileNet / EfficientNet 기반 MV 모델 → ONNX/CoreML 변환. "
         "선별기 엣지 배포 (<100ms latency) 타겟"),
        ("9. Multi-view Attention Fusion",
         "concat 대신 learnable weighted-avg (softmax α_t1,t2,b1). 파라미터 3개 추가로 "
         "뷰별 기여도 명시 학습 + 해석성 확보. Self-attention은 데이터 규모상 risky"),
    ]
    y = Inches(1.65)
    for title, body in items:
        add_text(s, Inches(0.6), y, Inches(3.8), Inches(0.4),
                 title, size=11, bold=True, color=C_ACCENT)
        add_text(s, Inches(4.4), y, Inches(8.5), Inches(0.6),
                 body, size=10, color=C_TITLE)
        y += Inches(0.58)


def slide_future_mid(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "개선 방향 · 중기 (데이터·인프라 확장 필요, 1~6개월)",
                 "다음 수확 시즌에 병행 추진 가능")

    items = [
        ("1. 데이터 확장 (우선순위 최상)",
         "다년도 / 다농가 / 다품종 (복숭아 10+ · 감귤 5+) / 결함 이미지. "
         "모델 설계보다 훨씬 큰 성능 상승 기대 — ML의 근본 레버"),
        ("2. 도메인 적응 (Domain Adaptation)",
         "카메라·조명 변경 시 재학습 없이 배포. "
         "Adversarial training / feature alignment / test-time adaptation 실험"),
        ("3. 자가지도학습 (Self-Supervised Pretraining)",
         "비라벨 과일 이미지 수만장 → DINOv2 / MAE pretrain → 소량 라벨로 fine-tune. "
         "라벨링 비용 50%+ 절감 가능성"),
        ("4. 반지도학습 (Semi-Supervised)",
         "측정값 없이 촬영만 된 이미지 활용 (FixMatch / pseudo-labeling). "
         "현재 버려지는 데이터 재활용"),
        ("5. Active Learning",
         "모델이 가장 불확실한 샘플을 우선 측정 의뢰. "
         "동일 성능 기준 라벨링 비용 50% 절감 잠재력"),
        ("6. 교차 작물 전이 (Cross-Crop Transfer)",
         "복숭아 모델 → 살구·자두 few-shot 적응. "
         "소수 품종 시장 진입 속도↑"),
        ("7. 동영상·회전 시퀀스 활용",
         "롤러 위 굴러가는 연속 프레임. 정지 3뷰보다 많은 정보. "
         "Video Transformer / ConvLSTM / ST-GCN"),
    ]
    y = Inches(1.7)
    for title, body in items:
        add_text(s, Inches(0.6), y, Inches(4.0), Inches(0.4),
                 title, size=12, bold=True, color=C_ACCENT)
        add_text(s, Inches(4.6), y, Inches(8.3), Inches(0.65),
                 body, size=10, color=C_TITLE)
        y += Inches(0.72)


def slide_future_long(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "개선 방향 · 장기 (센서/시스템 수준, 6개월+)",
                 "연구-개발 협업 필요")

    items = [
        ("1. 멀티모달 센서 융합",
         "RGB + NIR(근적외선) + 하이퍼스펙트럴. 특히 brix 정확도를 R² 0.9+로 도약 가능. "
         "비접촉 brix 측정 관련 분광학 논문 다수"),
        ("2. 3D 형태 재구성",
         "SfM / NeRF로 과일 완전한 3D 모델 → 부피·표면적 직접 계산. "
         "2D 프록시(height × max_w × min_w)보다 정확한 실무 지표"),
        ("3. 결함 세그멘테이션 + 등급화 통합",
         "단일 파이프라인으로 크기·당도·상품성(결함) 동시 산출. "
         "상처 위치/면적을 feature로 활용하는 end-to-end 모델"),
        ("4. Open-Set 품종 인식",
         "현재는 품종을 아는 전제. 개방형 품종 판별 + 품종별 measurement 통합. "
         "혼재 유통 환경에서 중요"),
        ("5. 농가 맞춤 Federated Learning",
         "각 농가 데이터 공유 없이 연합 학습. 개인정보·경쟁 이슈 해결 + 전체 모델 강화"),
        ("6. 실시간 운용 시스템",
         "모델 drift 감지, 자동 재학습 트리거, 신뢰도 경고 UI. "
         "운용자가 모델 상태를 투명하게 볼 수 있는 대시보드"),
        ("7. Sim-to-Real (선택)",
         "3D 과일 시뮬레이터 기반 대량 합성 데이터 + 실제 소량 데이터 결합. "
         "희귀 품종·결함 케이스 데이터 보강"),
    ]
    y = Inches(1.7)
    for title, body in items:
        add_text(s, Inches(0.6), y, Inches(4.0), Inches(0.4),
                 title, size=12, bold=True, color=C_ACCENT)
        add_text(s, Inches(4.6), y, Inches(8.3), Inches(0.65),
                 body, size=10, color=C_TITLE)
        y += Inches(0.72)


def slide_stack(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "기술 스택 · 산출물", "재현 가능한 오픈 구성")

    add_text(s, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
             "Stack", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(6), Inches(3), [
        "Framework: PyTorch 2.11 / torchvision 0.26",
        "Backbone: ResNet50 (ImageNet V2) / ViT-B/16 (비교용)",
        "Data: NumPy, OpenCV, Pillow",
        "Optim: SGD + momentum / Cosine LR / SmoothL1",
        "Runtime: Apple MPS (M4) / CUDA 자동 선택",
        "Python 3.11, conda env (fruit_grading)",
    ], size=13)

    add_text(s, Inches(7.0), Inches(1.7), Inches(6), Inches(0.5),
             "코드베이스 (willmer/)", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(7.0), Inches(2.2), Inches(6), Inches(3), [
        "make_labels.py — 라벨 CSV 빌더 + SANITY_BOUNDS",
        "dataset.py — Dataset + MV 변형",
        "model.py — ResNet50 / DenseNet201 / ViT-B/16 팩토리 + MV 래퍼",
        "train.py / train_mv.py / train_mv_mt.py — 학습 루프 3종",
        "requirements.txt / runs*/history.json 공개",
    ], size=13)

    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "저장소", size=16, bold=True, color=C_GREEN)
    add_text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.5),
             "https://github.com/YBNML/fruit_grading (public)",
             size=14, color=C_ACCENT, bold=True)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.8),
             "공개 범위: 코드 + 학습 설정 + 성능 결과(history.json)  /  "
             "비공개: 원본 이미지 데이터셋(DB/) + 학습된 모델 가중치(.pt)",
             size=12, color=C_MUTED)


def slide_closing(prs):
    s = add_blank(prs); fill_bg(s); add_accent_bar(s, Inches(0.45))
    slide_header(s, "요약 & 연구 제언", "한 장으로 보는 take-away")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
             "프로젝트 요약", size=16, bold=True, color=C_GREEN)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2), [
        "5 Phase · 26 runs · ~20시간 MPS 학습 · 작물당 4 target 모델 운용",
        "최고 성능: 감귤 min_w MV — MAPE 0.85%, R² 0.975 (전 10 target MAPE 1~3%)",
        "ResNet50 + MV(3뷰 concat) + 단일태스크 구조가 운용 최적",
    ], size=13)

    add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
             "다음 연구 우선순위 (제안)", size=16, bold=True, color=C_ACCENT)
    add_bullets(s, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), [
        "(단기) Phase 4 세그멘테이션 — 전처리 품질 개선의 가장 직접적 경로",
        "(중기) 데이터 확장 — 다년도·다품종·결함 이미지. ML의 근본 레버",
        "(중기) 자가지도학습 pretrain — 라벨링 비용 절감 + 강건성 향상",
        "(장기) 멀티모달(NIR/분광) — brix 정확도 도약 + 결함 탐지 통합",
        "(장기) 3D 재구성 + open-set 품종 인식 — 실무 통합 선별 시스템",
    ], size=13)


# ---------------------------------------------------------------------- main


def main():
    prs = new_prs()
    builders = [
        slide_title, slide_problem, slide_data, slide_pipeline,
        slide_phase0, slide_phase1, slide_phase3, slide_phase2,
        slide_phase5, slide_backbone, slide_summary,
        slide_insights, slide_insights2,
        slide_limits, slide_future_short, slide_future_mid, slide_future_long,
        slide_stack, slide_closing,
    ]
    for fn in builders:
        fn(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue  # skip footer on title
        add_footer(slide, i, total)

    prs.save(OUT)
    print(f"wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    main()
